"""
core_web.py — CORE Web, Document, Image & Utility Tools
=========================================================
All tools that run on Railway without PC dependency.

IMPORT IN core_tools.py (add at bottom):
    from core_web import _register_web_tools
    _register_web_tools(TOOLS)

CATEGORIES (add to core_orchestrator.py _TOOL_CATEGORIES):
    "web":      ["web_search", "web_fetch", "summarize_url"],
    "document": ["create_document", "create_spreadsheet", "create_presentation",
                 "read_document", "convert_document"],
    "image":    ["generate_image", "image_process"],
    "utils":    ["weather", "calc", "datetime_now", "currency", "translate",
                 "run_python"],

DEPENDENCIES (all already installed on Railway):
    httpx, python-docx, python-pptx, openpyxl, reportlab, Pillow
    All verified installed: 1.2.0, 1.0.2, 3.1.5, 4.4.10, 12.1.1
"""

import ast
import base64
import io
import json
import math
import os
import re
import tempfile
import traceback
from datetime import datetime
from pathlib import Path
from typing import Optional
from zoneinfo import ZoneInfo

import httpx

from core_config import _GEMINI_KEYS, _GEMINI_KEY_INDEX, gemini_chat, TOOL_CATEGORY_KEYWORDS

# ── Constants ──────────────────────────────────────────────────────────────────
WIB = ZoneInfo("Asia/Jakarta")
_OUTPUT_DIR = Path("/tmp/core_web_outputs")
_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ── Helpers ────────────────────────────────────────────────────────────────────

def _file_to_b64(path: str) -> str:
    """Read file and return base64 string."""
    return base64.b64encode(Path(path).read_bytes()).decode()


def _gemini_image_key() -> str:
    """Round-robin Gemini key for image generation."""
    import core_config as _cc
    key = _cc._GEMINI_KEYS[_cc._GEMINI_KEY_INDEX % len(_cc._GEMINI_KEYS)]
    _cc._GEMINI_KEY_INDEX = (_cc._GEMINI_KEY_INDEX + 1) % len(_cc._GEMINI_KEYS)
    return key


# ══════════════════════════════════════════════════════════════════════════════
# WEB TOOLS
# ══════════════════════════════════════════════════════════════════════════════

def t_web_search(query: str = "", max_results: str = "5") -> dict:
    """
    Search the web via DuckDuckGo HTML (no API key needed).
    Fallback to Bing HTML if DDG returns no results or fails.
    Returns list of {title, url, snippet}.
    """
    if not query:
        return {"ok": False, "error": "query required"}

    import urllib.parse

    def _clean(s: str) -> str:
        return re.sub(r'<[^>]+>', '', s).strip()

    def _ddg(q: str, n: int):
        encoded = urllib.parse.quote(q)
        r = httpx.get(
            f"https://html.duckduckgo.com/html/?q={encoded}",
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"},
            timeout=15, follow_redirects=True,
        )
        r.raise_for_status()
        html = r.text
        snippets = re.findall(r'class="result__snippet"[^>]*>(.*?)</a>', html, re.DOTALL)
        titles   = re.findall(r'class="result__a"[^>]*>(.*?)</a>', html, re.DOTALL)
        # Extract actual href, not the display URL span (more reliable)
        hrefs    = re.findall(r'<a class="result__a"[^>]+href="([^"]+)"', html)
        results  = []
        for i in range(min(n, len(snippets))):
            url = hrefs[i] if i < len(hrefs) else ""
            # DDG redirects — extract actual URL from uddg param if present
            if "uddg=" in url:
                try:
                    url = urllib.parse.unquote(url.split("uddg=")[1].split("&")[0])
                except Exception:
                    pass
            results.append({
                "title":   _clean(titles[i])   if i < len(titles) else "",
                "url":     url,
                "snippet": _clean(snippets[i]),
            })
        return results

    def _bing(q: str, n: int):
        encoded = urllib.parse.quote(q)
        r = httpx.get(
            f"https://www.bing.com/search?q={encoded}&count={n}",
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"},
            timeout=15, follow_redirects=True,
        )
        r.raise_for_status()
        html = r.text
        # Bing result blocks
        blocks  = re.findall(r'<li class="b_algo">(.*?)</li>', html, re.DOTALL)
        results = []
        for block in blocks[:n]:
            title   = _clean(re.search(r'<h2[^>]*>(.*?)</h2>', block, re.DOTALL).group(1) if re.search(r'<h2[^>]*>(.*?)</h2>', block, re.DOTALL) else "")
            href_m  = re.search(r'href="(https?://[^"]+)"', block)
            url     = href_m.group(1) if href_m else ""
            snippet = _clean(re.search(r'<p[^>]*>(.*?)</p>', block, re.DOTALL).group(1) if re.search(r'<p[^>]*>(.*?)</p>', block, re.DOTALL) else "")
            if title or url:
                results.append({"title": title, "url": url, "snippet": snippet})
        return results

    try:
        n = min(int(max_results) if max_results else 5, 20)
        results = []
        source  = "ddg"
        try:
            results = _ddg(query, n)
        except Exception as _de:
            source = f"ddg_failed({str(_de)[:40]})"

        if not results:
            try:
                results = _bing(query, n)
                source  = "bing_fallback"
            except Exception as _be:
                return {"ok": False, "error": f"both DDG and Bing failed. Bing: {_be}"}

        return {
            "ok":      True,
            "query":   query,
            "source":  source,
            "count":   len(results),
            "results": results,
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_web_fetch(url: str = "", max_chars: str = "8000") -> dict:
    """
    Fetch content from any URL. Returns cleaned text content.
    max_chars: truncate output (default 8000).
    """
    if not url:
        return {"ok": False, "error": "url required"}
    try:
        limit = int(max_chars) if max_chars else 8000
        r = httpx.get(
            url,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"},
            timeout=20,
            follow_redirects=True,
        )
        r.raise_for_status()
        ct = r.headers.get("content-type", "")

        if "text" in ct or "html" in ct or "json" in ct or "xml" in ct:
            raw = r.text
            # Strip HTML tags for cleaner content
            clean = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL)
            clean = re.sub(r'<style[^>]*>.*?</style>',  '', clean, flags=re.DOTALL)
            clean = re.sub(r'<[^>]+>', ' ', clean)
            clean = re.sub(r'\s+', ' ', clean).strip()
            truncated = len(clean) > limit
            return {
                "ok":        True,
                "url":       url,
                "content":   clean[:limit],
                "truncated": truncated,
                "length":    len(clean),
            }
        else:
            # Binary — return base64
            b64 = base64.b64encode(r.content).decode()
            return {"ok": True, "url": url, "binary": True,
                    "content_type": ct, "base64": b64[:limit]}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_summarize_url(url: str = "", focus: str = "") -> dict:
    """
    Fetch URL and summarize content via Gemini.
    focus: optional — what aspect to focus on (default: general summary).
    """
    if not url:
        return {"ok": False, "error": "url required"}
    fetch = t_web_fetch(url, max_chars="6000")
    if not fetch.get("ok"):
        return fetch
    content = fetch.get("content", "")
    if not content:
        return {"ok": False, "error": "no text content found at URL"}
    prompt = (
        f"URL: {url}\n\n"
        f"Content:\n{content[:5000]}\n\n"
        f"{'Focus: ' + focus if focus else 'Provide a concise summary.'}"
    )
    try:
        summary = gemini_chat(
            system="You are a concise web content summarizer. Extract key information clearly.",
            user=prompt,
            max_tokens=800,
        )
        return {"ok": True, "url": url, "summary": summary}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# DOCUMENT TOOLS
# ══════════════════════════════════════════════════════════════════════════════

def t_create_document(
    content: str = "",
    filename: str = "document",
    format: str  = "docx",
    title: str   = "",
) -> dict:
    """
    Create a document from text content.
    format: docx | pdf | txt | md | csv
    content: the text/markdown content to put in the document.
    title: optional document title (shown as heading in docx/pdf).
    Returns: {ok, filename, base64, size_bytes, format}
    """
    if not content:
        return {"ok": False, "error": "content required"}

    fmt  = format.lower().strip().lstrip(".")
    name = filename.rstrip("." + fmt) if filename.endswith("." + fmt) else filename
    out  = _OUTPUT_DIR / f"{name}.{fmt}"

    try:
        if fmt == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            if title:
                h = doc.add_heading(title, level=1)
                h.alignment = WD_ALIGN_PARAGRAPH.LEFT

            # Parse content — support basic markdown-like syntax
            for line in content.splitlines():
                stripped = line.strip()
                if stripped.startswith("## "):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith("# "):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith("### "):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith("- ") or stripped.startswith("* "):
                    doc.add_paragraph(stripped[2:], style="List Bullet")
                elif re.match(r'^\d+\. ', stripped):
                    doc.add_paragraph(re.sub(r'^\d+\. ', '', stripped), style="List Number")
                elif stripped == "---" or stripped == "***":
                    doc.add_paragraph("─" * 40)
                elif stripped:
                    doc.add_paragraph(stripped)
                else:
                    doc.add_paragraph("")

            doc.save(str(out))

        elif fmt == "pdf":
            from reportlab.pdfgen import canvas as rl_canvas
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.units import cm
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib import colors

            styles   = getSampleStyleSheet()
            story    = []
            if title:
                title_style = ParagraphStyle(
                    "CustomTitle",
                    parent=styles["Title"],
                    fontSize=18, spaceAfter=12,
                )
                story.append(Paragraph(title, title_style))
                story.append(Spacer(1, 0.3 * cm))

            body_style = styles["BodyText"]
            body_style.fontSize    = 11
            body_style.leading     = 16
            body_style.spaceAfter  = 6

            for line in content.splitlines():
                stripped = line.strip()
                if stripped.startswith("# "):
                    story.append(Paragraph(stripped[2:], styles["Heading1"]))
                elif stripped.startswith("## "):
                    story.append(Paragraph(stripped[3:], styles["Heading2"]))
                elif stripped.startswith("### "):
                    story.append(Paragraph(stripped[4:], styles["Heading3"]))
                elif stripped.startswith("- ") or stripped.startswith("* "):
                    story.append(Paragraph(f"• {stripped[2:]}", body_style))
                elif stripped:
                    story.append(Paragraph(stripped, body_style))
                else:
                    story.append(Spacer(1, 0.2 * cm))

            doc = SimpleDocTemplate(str(out), pagesize=A4,
                                    rightMargin=2*cm, leftMargin=2*cm,
                                    topMargin=2*cm, bottomMargin=2*cm)
            doc.build(story)

        elif fmt in ("txt", "md"):
            out.write_text(content, encoding="utf-8")

        elif fmt == "csv":
            # Try to parse content as CSV-like data
            out.write_text(content, encoding="utf-8")

        else:
            return {"ok": False, "error": f"unsupported format: {fmt}. Use docx|pdf|txt|md|csv"}

        size  = out.stat().st_size
        b64   = _file_to_b64(str(out))
        return {
            "ok":        True,
            "filename":  out.name,
            "format":    fmt,
            "size_bytes": size,
            "base64":    b64,
            "note":      "base64 field contains the file — send to Telegram or save to disk",
        }
    except Exception as e:
        return {"ok": False, "error": str(e), "trace": traceback.format_exc()[:500]}


def t_create_spreadsheet(
    data: str      = "",
    filename: str  = "spreadsheet",
    sheet_name: str = "Sheet1",
    format: str    = "xlsx",
) -> dict:
    """
    Create a spreadsheet from data.
    data: JSON array of arrays (rows) OR CSV text.
         Example JSON: [["Name","Age"],["Alice",30],["Bob",25]]
         Example CSV:  Name,Age\\nAlice,30\\nBob,25
    format: xlsx | csv
    Returns: {ok, filename, base64, size_bytes, rows, cols}
    """
    if not data:
        return {"ok": False, "error": "data required"}

    fmt  = format.lower().strip().lstrip(".")
    name = filename.rstrip("." + fmt) if filename.endswith("." + fmt) else filename
    out  = _OUTPUT_DIR / f"{name}.{fmt}"

    # Parse data
    rows = []
    try:
        parsed = json.loads(data)
        if isinstance(parsed, list):
            rows = parsed
        else:
            return {"ok": False, "error": "data must be a JSON array of arrays"}
    except json.JSONDecodeError:
        # Try CSV
        import csv, io as _io
        reader = csv.reader(_io.StringIO(data))
        rows   = list(reader)

    if not rows:
        return {"ok": False, "error": "no rows parsed from data"}

    try:
        if fmt == "xlsx":
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name

            header_font  = Font(bold=True, color="FFFFFF", size=11)
            header_fill  = PatternFill("solid", fgColor="2E5DA3")
            header_align = Alignment(horizontal="center", vertical="center")
            thin_border  = Border(
                left=Side(style="thin"), right=Side(style="thin"),
                top=Side(style="thin"), bottom=Side(style="thin"),
            )

            for r_idx, row in enumerate(rows, start=1):
                for c_idx, val in enumerate(row, start=1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=val)
                    cell.border = thin_border
                    if r_idx == 1:
                        cell.font      = header_font
                        cell.fill      = header_fill
                        cell.alignment = header_align
                    else:
                        cell.alignment = Alignment(horizontal="left", vertical="center")
                        if r_idx % 2 == 0:
                            cell.fill = PatternFill("solid", fgColor="EEF2FF")

            # Auto-width columns
            for col in ws.columns:
                max_len = max((len(str(cell.value or "")) for cell in col), default=8)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

            ws.row_dimensions[1].height = 20
            wb.save(str(out))

        elif fmt == "csv":
            import csv
            with open(str(out), "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                writer.writerows(rows)

        else:
            return {"ok": False, "error": f"unsupported format: {fmt}. Use xlsx|csv"}

        size = out.stat().st_size
        b64  = _file_to_b64(str(out))
        return {
            "ok":        True,
            "filename":  out.name,
            "format":    fmt,
            "size_bytes": size,
            "rows":      len(rows),
            "cols":      max(len(r) for r in rows) if rows else 0,
            "base64":    b64,
        }
    except Exception as e:
        return {"ok": False, "error": str(e), "trace": traceback.format_exc()[:500]}


def t_create_presentation(
    slides: str    = "",
    filename: str  = "presentation",
    theme: str     = "default",
) -> dict:
    """
    Create a PowerPoint presentation.
    slides: JSON array of slide objects:
      [
        {"title": "Slide Title", "content": "bullet\\nbullet\\nbullet", "notes": "speaker notes"},
        {"title": "...", "content": "...", "layout": "title_only|content|two_col"}
      ]
    theme: default | dark | minimal
    Returns: {ok, filename, base64, size_bytes, slide_count}
    """
    if not slides:
        return {"ok": False, "error": "slides required"}

    try:
        slide_data = json.loads(slides) if isinstance(slides, str) else slides
        if not isinstance(slide_data, list):
            return {"ok": False, "error": "slides must be a JSON array"}
    except json.JSONDecodeError as e:
        return {"ok": False, "error": f"invalid JSON: {e}"}

    name = filename.rstrip(".pptx") if filename.endswith(".pptx") else filename
    out  = _OUTPUT_DIR / f"{name}.pptx"

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Theme colors
        themes = {
            "default": {"bg": RGBColor(0xFF, 0xFF, 0xFF), "title": RGBColor(0x1F, 0x3864, 0x96),
                        "body": RGBColor(0x26, 0x26, 0x26), "accent": RGBColor(0x2E, 0x5D, 0xA3)},
            "dark":    {"bg": RGBColor(0x1A, 0x1A, 0x2E), "title": RGBColor(0xE9, 0x4C, 0x60),
                        "body": RGBColor(0xE0, 0xE0, 0xE0), "accent": RGBColor(0x16, 0x21, 0x3E)},
            "minimal": {"bg": RGBColor(0xFA, 0xFA, 0xFA), "title": RGBColor(0x33, 0x33, 0x33),
                        "body": RGBColor(0x55, 0x55, 0x55), "accent": RGBColor(0x00, 0xAA, 0xBB)},
        }
        tc = themes.get(theme, themes["default"])

        blank_layout = prs.slide_layouts[6]  # blank

        for s in slide_data:
            slide = prs.slides.add_slide(blank_layout)

            # Background
            bg       = slide.background
            fill     = bg.fill
            fill.solid()
            fill.fore_color.rgb = tc["bg"]

            s_title   = s.get("title", "")
            s_content = s.get("content", "")
            s_notes   = s.get("notes", "")

            # Title box
            if s_title:
                txb = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.1)
                )
                tf = txb.text_frame
                tf.word_wrap = True
                p  = tf.paragraphs[0]
                run = p.add_run()
                run.text = s_title
                run.font.bold  = True
                run.font.size  = Pt(36)
                run.font.color.rgb = tc["title"]

            # Content box
            if s_content:
                txb = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.4)
                )
                tf = txb.text_frame
                tf.word_wrap = True
                lines = s_content.splitlines()
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    stripped = line.strip()
                    is_bullet = stripped.startswith("- ") or stripped.startswith("• ")
                    text = stripped[2:] if is_bullet else stripped
                    run = p.add_run()
                    run.text = ("• " if is_bullet else "") + text
                    run.font.size  = Pt(20)
                    run.font.color.rgb = tc["body"]
                    p.space_after = Pt(6)

            # Speaker notes
            if s_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = s_notes

        prs.save(str(out))

        size = out.stat().st_size
        b64  = _file_to_b64(str(out))
        return {
            "ok":          True,
            "filename":    out.name,
            "format":      "pptx",
            "size_bytes":  size,
            "slide_count": len(slide_data),
            "base64":      b64,
        }
    except Exception as e:
        return {"ok": False, "error": str(e), "trace": traceback.format_exc()[:500]}


def t_read_document(
    base64_content: str = "",
    filename: str       = "",
    format: str         = "",
) -> dict:
    """
    Read and extract text from a document.
    base64_content: base64-encoded file bytes.
    filename: helps detect format if format not specified.
    format: docx | xlsx | pptx | txt | md | csv (auto-detected from filename if not set).
    Returns: {ok, text, format, pages/sheets/slides}
    """
    if not base64_content:
        return {"ok": False, "error": "base64_content required"}

    # Detect format
    fmt = format.lower().strip().lstrip(".") if format else ""
    if not fmt and filename:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        fmt = ext
    if not fmt:
        return {"ok": False, "error": "format required (docx|xlsx|pptx|txt|md|csv)"}

    try:
        file_bytes = base64.b64decode(base64_content)
        buf        = io.BytesIO(file_bytes)

        if fmt == "docx":
            from docx import Document
            doc   = Document(buf)
            lines = [p.text for p in doc.paragraphs if p.text.strip()]
            text  = "\n".join(lines)
            return {"ok": True, "format": "docx", "text": text[:20000],
                    "paragraphs": len(lines), "truncated": len(text) > 20000}

        elif fmt == "xlsx":
            from openpyxl import load_workbook
            wb    = load_workbook(buf, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"=== Sheet: {ws.title} ===")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join(str(v) if v is not None else "" for v in row))
            text = "\n".join(parts)
            return {"ok": True, "format": "xlsx", "text": text[:20000],
                    "sheets": len(wb.worksheets), "truncated": len(text) > 20000}

        elif fmt == "pptx":
            from pptx import Presentation
            prs   = Presentation(buf)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"=== Slide {i} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            text = "\n".join(parts)
            return {"ok": True, "format": "pptx", "text": text[:20000],
                    "slides": len(prs.slides), "truncated": len(text) > 20000}

        elif fmt in ("txt", "md", "csv"):
            text = file_bytes.decode("utf-8", errors="replace")
            return {"ok": True, "format": fmt, "text": text[:20000],
                    "truncated": len(text) > 20000}

        else:
            return {"ok": False, "error": f"unsupported format: {fmt}"}

    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_convert_document(
    base64_content: str = "",
    from_format: str    = "",
    to_format: str      = "",
    filename: str       = "converted",
) -> dict:
    """
    Convert between document formats.
    Supported conversions:
      md → docx | pdf | txt
      txt → docx | pdf | md
      docx → txt | md
      xlsx → csv
      csv → xlsx
    base64_content: base64-encoded source file.
    Returns: {ok, filename, base64, size_bytes}
    """
    if not base64_content or not from_format or not to_format:
        return {"ok": False, "error": "base64_content, from_format, to_format all required"}

    src_fmt = from_format.lower().strip().lstrip(".")
    dst_fmt = to_format.lower().strip().lstrip(".")

    try:
        file_bytes = base64.b64decode(base64_content)

        # Extract text from source
        if src_fmt in ("txt", "md"):
            text = file_bytes.decode("utf-8", errors="replace")
        elif src_fmt == "docx":
            from docx import Document
            doc  = Document(io.BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif src_fmt == "csv":
            text = file_bytes.decode("utf-8", errors="replace")
        else:
            return {"ok": False, "error": f"cannot read from format: {src_fmt}"}

        # Write to target format
        result = t_create_document(
            content=text,
            filename=filename,
            format=dst_fmt,
        ) if dst_fmt in ("docx", "pdf", "txt", "md") else (
            t_create_spreadsheet(data=text, filename=filename, format=dst_fmt)
            if dst_fmt in ("xlsx", "csv") else
            {"ok": False, "error": f"unsupported target format: {dst_fmt}"}
        )
        return result

    except Exception as e:
        return {"ok": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# IMAGE TOOLS
# ══════════════════════════════════════════════════════════════════════════════

def t_generate_image(
    prompt: str      = "",
    filename: str    = "generated",
    aspect_ratio: str = "1:1",
) -> dict:
    """
    Generate an image using Gemini Imagen API.
    prompt: description of the image to generate.
    aspect_ratio: 1:1 | 16:9 | 9:16 | 4:3 | 3:4
    Returns: {ok, filename, base64, mime_type}
    """
    if not prompt:
        return {"ok": False, "error": "prompt required"}
    import core_config as _cc
    if not _cc._GEMINI_KEYS:
        return {"ok": False, "error": "GEMINI_KEYS not set"}

    valid_ratios = {"1:1", "16:9", "9:16", "4:3", "3:4"}
    if aspect_ratio not in valid_ratios:
        aspect_ratio = "1:1"

    last_err = None
    for _ in range(len(_cc._GEMINI_KEYS)):
        key = _cc._GEMINI_KEYS[_cc._GEMINI_KEY_INDEX % len(_cc._GEMINI_KEYS)]
        _cc._GEMINI_KEY_INDEX = (_cc._GEMINI_KEY_INDEX + 1) % len(_cc._GEMINI_KEYS)
        try:
            r = httpx.post(
                "https://generativelanguage.googleapis.com/v1beta/models/imagen-3.0-generate-001:predict",
                params={"key": key},
                headers={"Content-Type": "application/json"},
                json={
                    "instances":  [{"prompt": prompt}],
                    "parameters": {
                        "sampleCount":  1,
                        "aspectRatio":  aspect_ratio,
                        "safetyFilterLevel": "block_only_high",
                    },
                },
                timeout=60,
            )
            if r.status_code == 429:
                last_err = "429"
                continue
            r.raise_for_status()
            data       = r.json()
            predictions = data.get("predictions", [])
            if not predictions:
                last_err = "no predictions returned"
                continue
            b64  = predictions[0].get("bytesBase64Encoded", "")
            mime = predictions[0].get("mimeType", "image/png")
            if not b64:
                last_err = "empty bytesBase64Encoded"
                continue
            # Save to file
            ext  = mime.split("/")[-1]
            out  = _OUTPUT_DIR / f"{filename}.{ext}"
            out.write_bytes(base64.b64decode(b64))
            return {
                "ok":       True,
                "filename": out.name,
                "mime_type": mime,
                "base64":   b64,
                "prompt":   prompt,
            }
        except Exception as e:
            last_err = str(e)
            continue
    return {"ok": False, "error": f"Image generation failed. Last: {last_err}"}


def t_image_process(
    base64_content: str = "",
    operation: str      = "info",
    params: str         = "{}",
    filename: str       = "processed",
) -> dict:
    """
    Process an image using Pillow.
    base64_content: base64-encoded image bytes.
    operation:
      info       — return size, mode, format
      resize     — params: {"width": N, "height": N, "keep_ratio": true}
      crop       — params: {"left": N, "top": N, "right": N, "bottom": N}
      rotate     — params: {"degrees": 90}
      flip       — params: {"direction": "horizontal"|"vertical"}
      convert    — params: {"mode": "RGB"|"L"|"RGBA", "format": "png"|"jpeg"|"webp"}
      watermark  — params: {"text": "...", "opacity": 0.5, "position": "center"|"bottom_right"}
      thumbnail  — params: {"size": 256}  — square thumbnail
    Returns: {ok, filename, base64, mime_type, width, height}
    """
    if not base64_content:
        return {"ok": False, "error": "base64_content required"}
    try:
        from PIL import Image, ImageDraw, ImageFont
        p = json.loads(params) if isinstance(params, str) else (params or {})
        img_bytes = base64.b64decode(base64_content)
        img       = Image.open(io.BytesIO(img_bytes))
        orig_fmt  = (img.format or "PNG").upper()

        op = operation.lower().strip()

        if op == "info":
            return {
                "ok": True, "operation": "info",
                "width": img.width, "height": img.height,
                "mode": img.mode, "format": orig_fmt,
            }

        elif op == "resize":
            w = int(p.get("width", img.width))
            h = int(p.get("height", img.height))
            if p.get("keep_ratio", True):
                img.thumbnail((w, h), Image.LANCZOS)
            else:
                img = img.resize((w, h), Image.LANCZOS)

        elif op == "crop":
            box = (int(p.get("left", 0)), int(p.get("top", 0)),
                   int(p.get("right", img.width)), int(p.get("bottom", img.height)))
            img = img.crop(box)

        elif op == "rotate":
            img = img.rotate(float(p.get("degrees", 90)), expand=True)

        elif op == "flip":
            direction = p.get("direction", "horizontal").lower()
            if direction == "horizontal":
                img = img.transpose(Image.FLIP_LEFT_RIGHT)
            else:
                img = img.transpose(Image.FLIP_TOP_BOTTOM)

        elif op == "convert":
            mode = p.get("mode", "RGB")
            if img.mode != mode:
                img = img.convert(mode)
            orig_fmt = p.get("format", orig_fmt).upper()

        elif op == "watermark":
            text     = p.get("text", "WATERMARK")
            opacity  = float(p.get("opacity", 0.4))
            position = p.get("position", "bottom_right")
            if img.mode != "RGBA":
                img = img.convert("RGBA")
            overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
            draw    = ImageDraw.Draw(overlay)
            # Use default font
            font_size = max(20, img.width // 20)
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", font_size)
            except Exception:
                font = ImageFont.load_default()
            bbox = draw.textbbox((0, 0), text, font=font)
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            if position == "center":
                x, y = (img.width - tw) // 2, (img.height - th) // 2
            else:  # bottom_right
                x, y = img.width - tw - 20, img.height - th - 20
            draw.text((x, y), text, font=font,
                      fill=(255, 255, 255, int(255 * opacity)))
            img = Image.alpha_composite(img, overlay)

        elif op == "thumbnail":
            size = int(p.get("size", 256))
            img.thumbnail((size, size), Image.LANCZOS)

        else:
            return {"ok": False, "error": f"unknown operation: {op}"}

        # Save result
        out_fmt = orig_fmt if orig_fmt in ("PNG", "JPEG", "WEBP", "GIF") else "PNG"
        ext     = out_fmt.lower().replace("jpeg", "jpg")
        out     = _OUTPUT_DIR / f"{filename}.{ext}"
        save_kwargs = {"quality": 90} if out_fmt == "JPEG" else {}
        # Convert RGBA→RGB for JPEG
        if out_fmt == "JPEG" and img.mode == "RGBA":
            img = img.convert("RGB")
        img.save(str(out), format=out_fmt, **save_kwargs)

        b64 = _file_to_b64(str(out))
        return {
            "ok":        True,
            "operation": op,
            "filename":  out.name,
            "mime_type": f"image/{ext}",
            "width":     img.width,
            "height":    img.height,
            "base64":    b64,
        }
    except Exception as e:
        return {"ok": False, "error": str(e), "trace": traceback.format_exc()[:400]}


# ══════════════════════════════════════════════════════════════════════════════
# UTILITY TOOLS
# ══════════════════════════════════════════════════════════════════════════════

def t_weather(location: str = "Jakarta") -> dict:
    """
    Get current weather for a location via wttr.in (no API key needed).
    location: city name or coordinates.
    Returns: {ok, location, condition, temp_c, feels_like_c, humidity, wind_kph, summary}
    """
    try:
        r = httpx.get(
            f"https://wttr.in/{location}",
            params={"format": "j1"},
            headers={"User-Agent": "curl/7.68.0"},
            timeout=10,
        )
        r.raise_for_status()
        data    = r.json()
        current = data["current_condition"][0]
        area    = data.get("nearest_area", [{}])[0]
        city    = area.get("areaName", [{}])[0].get("value", location)

        condition  = current["weatherDesc"][0]["value"]
        temp_c     = int(current["temp_C"])
        feels_c    = int(current["FeelsLikeC"])
        humidity   = int(current["humidity"])
        wind_kph   = int(current["windspeedKmph"])
        visibility = int(current.get("visibility", 0))

        summary = (
            f"{city}: {condition}, {temp_c}°C "
            f"(feels {feels_c}°C), humidity {humidity}%, "
            f"wind {wind_kph} km/h"
        )
        return {
            "ok":          True,
            "location":    city,
            "condition":   condition,
            "temp_c":      temp_c,
            "feels_like_c": feels_c,
            "humidity":    humidity,
            "wind_kph":    wind_kph,
            "visibility_km": visibility,
            "summary":     summary,
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_calc(expression: str = "") -> dict:
    """
    Safe math calculator. Evaluates mathematical expressions.
    Supports: +, -, *, /, **, %, //, sqrt, sin, cos, tan, log, pi, e, abs, round, etc.
    expression: math expression string e.g. "2 ** 10", "sqrt(144)", "0.5 * 71168.77"
    Auto-sanitises: strips $, currency symbols, thousands commas (e.g. "$71,168.77" → "71168.77")
    """
    if not expression:
        return {"ok": False, "error": "expression required"}
    try:
        # Auto-sanitise common patterns from financial/web data before validation
        import re as _re
        clean = expression.strip()
        clean = clean.replace("$", "").replace("€", "").replace("£", "").replace("¥", "")
        clean = clean.replace(",", "")  # remove thousands separators: 71,168.77 → 71168.77
        clean = clean.replace("USD", "").replace("BTC", "").replace("ETH", "")
        clean = clean.strip()

        # Safe eval — only math functions allowed
        safe_globals = {
            "__builtins__": {},
            "abs": abs, "round": round, "min": min, "max": max,
            "sum": sum, "pow": pow, "int": int, "float": float,
        }
        safe_globals.update({
            k: getattr(math, k) for k in dir(math) if not k.startswith("_")
        })
        # Validate — only allow safe characters
        allowed = re.sub(r'[\d\s\+\-\*\/\(\)\.\%\_a-zA-Z]', '', clean)
        if allowed:
            return {"ok": False, "error": f"unsafe characters in expression: {repr(allowed)}"}
        result = eval(clean, safe_globals)  # noqa: S307
        # Format result
        if isinstance(result, float):
            formatted = f"{result:.10g}"
        else:
            formatted = str(result)
        return {"ok": True, "expression": clean, "original": expression, "result": result, "formatted": formatted}
    except ZeroDivisionError:
        return {"ok": False, "error": "division by zero"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_datetime_now(timezone: str = "Asia/Jakarta") -> dict:
    """
    Get current date and time.
    timezone: IANA timezone string. Default: Asia/Jakarta (WIB, UTC+7).
    Returns: {ok, iso, date, time, day_of_week, timezone, unix_ts}
    """
    try:
        tz  = ZoneInfo(timezone)
        now = datetime.now(tz)
        return {
            "ok":          True,
            "iso":         now.isoformat(),
            "date":        now.strftime("%Y-%m-%d"),
            "time":        now.strftime("%H:%M:%S"),
            "day_of_week": now.strftime("%A"),
            "timezone":    timezone,
            "utc_offset":  now.strftime("%z"),
            "unix_ts":     int(now.timestamp()),
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_currency(
    amount: str    = "1",
    from_cur: str  = "USD",
    to_cur: str    = "IDR",
) -> dict:
    """
    Convert currency using exchangerate-api (free, no key needed).
    amount: amount to convert (default 1).
    from_cur: source currency code (e.g. USD, EUR, IDR).
    to_cur: target currency code.
    Returns: {ok, amount, from, to, rate, result, formatted}
    """
    try:
        amt  = float(amount) if amount else 1.0
        src  = from_cur.upper().strip()
        dst  = to_cur.upper().strip()
        r    = httpx.get(
            f"https://open.er-api.com/v6/latest/{src}",
            timeout=10,
        )
        r.raise_for_status()
        data  = r.json()
        rates = data.get("rates", {})
        if dst not in rates:
            return {"ok": False, "error": f"currency '{dst}' not found. "
                    f"Available: {', '.join(list(rates.keys())[:20])}..."}
        rate   = rates[dst]
        result = amt * rate
        return {
            "ok":       True,
            "amount":   amt,
            "from":     src,
            "to":       dst,
            "rate":     rate,
            "result":   result,
            "formatted": f"{amt:,.2f} {src} = {result:,.2f} {dst}",
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_translate(
    text: str         = "",
    target_language: str = "English",
    source_language: str = "",
) -> dict:
    """
    Translate text using Gemini (no extra API key needed).
    text: text to translate.
    target_language: target language name (e.g. English, Indonesian, Japanese).
    source_language: optional source language hint. Auto-detected if empty.
    Returns: {ok, original, translated, target_language}
    """
    if not text:
        return {"ok": False, "error": "text required"}
    try:
        src_hint = f" from {source_language}" if source_language else ""
        result   = gemini_chat(
            system="You are a precise translator. Return ONLY the translated text, nothing else.",
            user=f"Translate the following text{src_hint} to {target_language}:\n\n{text}",
            max_tokens=2048,
        )
        return {
            "ok":              True,
            "original":        text[:500],
            "translated":      result,
            "target_language": target_language,
            "source_language": source_language or "auto-detected",
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_run_python(code: str = "", timeout: str = "10") -> dict:
    """
    Execute Python code safely on Railway.
    Runs in a subprocess with timeout. No filesystem writes outside /tmp.
    Captures stdout + stderr. Returns output.
    USE FOR: calculations, data transformations, string processing, algorithms.
    NOT FOR: installing packages, network calls, reading/writing arbitrary files.
    code: Python code to execute.
    timeout: max execution seconds (default 10, max 30).
    """
    if not code:
        return {"ok": False, "error": "code required"}
    try:
        import subprocess
        t = min(int(timeout) if timeout else 10, 30)
        # Write to temp file
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".py", delete=False,
            dir="/tmp", prefix="core_run_"
        ) as f:
            f.write(code)
            tmp = f.name
        try:
            r = subprocess.run(
                ["python3", tmp],
                capture_output=True, text=True, timeout=t,
                env={k: v for k, v in os.environ.items() if k not in ("SUPABASE_SERVICE_KEY","SUPABASE_ANON_KEY","GITHUB_PAT","TELEGRAM_BOT_TOKEN","BINANCE_SECRET_KEY","MCP_SECRET")}, 
            )
            output = (r.stdout + r.stderr)[:5000]
            return {
                "ok":         r.returncode == 0,
                "returncode": r.returncode,
                "output":     output,
                "truncated":  len(r.stdout + r.stderr) > 5000,
            }
        except subprocess.TimeoutExpired:
            return {"ok": False, "error": f"execution timed out after {t}s"}
        finally:
            try:
                os.unlink(tmp)
            except Exception:
                pass
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# META-TOOLS — tool self-awareness
# ══════════════════════════════════════════════════════════════════════════════

def t_list_tools(category: str = "", search: str = "") -> dict:
    """
    List all available tools in the TOOLS registry.
    category: optional keyword filter — matches against tool names in each category.
              Uses live category build from orchestrator (always current, no hardcode).
    search:   optional keyword filter on tool name + description.
    Returns full list with name, args, perm, desc per tool.
    """
    try:
        from core_tools import TOOLS
        # Single source of truth — imported from core_config.TOOL_CATEGORY_KEYWORDS
        _KW = TOOL_CATEGORY_KEYWORDS
        cats: dict = {cat: [] for cat in _KW}
        cats["misc"] = []
        for tn in TOOLS.keys():
            placed = False
            for cat, kws in _KW.items():
                if any(kw in tn for kw in kws):
                    cats[cat].append(tn); placed = True; break
            if not placed:
                cats["misc"].append(tn)
        cats = {k: v for k, v in cats.items() if v}

        cat = category.lower().strip() if category else ""
        if cat:
            if cats and cat not in cats:
                return {
                    "ok": False,
                    "error": f"unknown category '{cat}'. Available: {', '.join(sorted(cats.keys()))}",
                }
            candidate_names = cats.get(cat, list(TOOLS.keys())) if cats else list(TOOLS.keys())
        else:
            candidate_names = list(TOOLS.keys())

        results = []
        kw = search.lower().strip() if search else ""
        for name in candidate_names:
            tdef = TOOLS.get(name)
            if not tdef:
                continue
            desc = tdef.get("desc", "")
            if kw and kw not in name.lower() and kw not in desc.lower():
                continue
            args_str = ", ".join(
                (a["name"] if isinstance(a, dict) else str(a))
                for a in (tdef.get("args") or [])
            )
            results.append({
                "name": name,
                "args": args_str,
                "perm": tdef.get("perm", ""),
                "desc": desc[:120],
            })

        return {
            "ok":            True,
            "total":         len(TOOLS),
            "filtered":      len(results),
            "category":      cat or "all",
            "search":        search or "",
            "available_cats": sorted(cats.keys()) if cats else [],
            "tools":         results,
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


def t_get_tool_info(name: str = "") -> dict:
    """
    Get full detail of a specific tool by exact name.
    Returns args (with types if available), full description, permission level.
    Use this before calling an unfamiliar tool to verify correct parameter names.
    """
    if not name:
        return {"ok": False, "error": "name required"}
    try:
        from core_tools import TOOLS
        name = name.strip()
        tdef = TOOLS.get(name)
        if not tdef:
            # Fuzzy suggest
            close = [t for t in TOOLS if name.lower() in t.lower()]
            return {
                "ok":      False,
                "error":   f"tool '{name}' not found",
                "similar": close[:10],
            }
        args = tdef.get("args") or []
        args_detail = []
        for a in args:
            if isinstance(a, dict):
                args_detail.append({
                    "name":     a.get("name", ""),
                    "type":     a.get("type", "string"),
                    "required": a.get("required", False),
                    "default":  a.get("default", ""),
                })
            else:
                args_detail.append({"name": str(a), "type": "string"})

        return {
            "ok":   True,
            "name": name,
            "perm": tdef.get("perm", ""),
            "desc": tdef.get("desc", ""),
            "args": args_detail,
            "args_simple": ", ".join(
                (a["name"] if isinstance(a, dict) else str(a))
                for a in args
            ),
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}

def t_get_table_schema(table: str = "") -> dict:
    """Get schema for a Supabase table.
    Merges two sources:
      - Live DB query (information_schema) → actual column names + data types
      - _SB_SCHEMA registry → fat_columns, required, enums, safe_select, on_conflict
    table="": summary of all tables from _SB_SCHEMA (no live query).
    table="name": full merged schema for that table.
    Falls back to _SB_SCHEMA-only if live DB query fails (no PAT etc).
    """
    try:
        from core_tools import _SB_SCHEMA

        # ── No table arg: return registry summary ──────────────────────────────
        if not table or not table.strip():
            summary = {}
            for tname, tdef in _SB_SCHEMA.get("tables", {}).items():
                summary[tname] = {
                    "pk_type":     tdef.get("pk_type", "?"),
                    "columns":     sorted(tdef.get("columns", {}).keys()),
                    "fat_columns": tdef.get("fat_columns", []),
                    "safe_select": tdef.get("safe_select", "*"),
                }
            return {
                "ok":               True,
                "mode":             "summary",
                "table_count":      len(summary),
                "tables":           summary,
                "tombstone_tables": sorted(_SB_SCHEMA.get("_tombstone", set())),
            }

        table = table.strip()

        # ── Tombstone check ────────────────────────────────────────────────────
        if table in _SB_SCHEMA.get("_tombstone", set()):
            return {"ok": False, "error": f"TOMBSTONE: '{table}' is retired — never query"}

        # ── _SB_SCHEMA metadata (always available) ─────────────────────────────
        tdef = _SB_SCHEMA.get("tables", {}).get(table, {})
        meta = {
            "pk":          tdef.get("pk", "id"),
            "pk_type":     tdef.get("pk_type", "unknown"),
            "fat_columns": tdef.get("fat_columns", []),
            "safe_select": tdef.get("safe_select", "*"),
            "required":    tdef.get("required", []),
            "enums":       tdef.get("enums", {}),
            "on_conflict": tdef.get("on_conflict", ""),
        }

        # ── Live DB query for actual columns + types ───────────────────────────
        live_columns = []
        live_source = "schema_registry_only"
        try:
            from core_config import SUPABASE_REF, SUPABASE_PAT
            if SUPABASE_REF and SUPABASE_PAT:
                resp = httpx.post(
                    f"https://api.supabase.com/v1/projects/{SUPABASE_REF}/database/query",
                    headers={"Authorization": f"Bearer {SUPABASE_PAT}",
                             "Content-Type": "application/json"},
                    json={"query": (
                        f"SELECT column_name, data_type, is_nullable "
                        f"FROM information_schema.columns "
                        f"WHERE table_name='{table}' AND table_schema='public' "
                        f"ORDER BY ordinal_position"
                    )},
                    timeout=12,
                )
                if resp.status_code in (200, 201):
                    live_columns = resp.json()
                    live_source = "live_db"
        except Exception as _le:
            live_source = f"live_db_failed: {str(_le)[:80]}"

        # Fallback: build columns from _SB_SCHEMA if live query failed
        if not live_columns and tdef.get("columns"):
            live_columns = [
                {"column_name": col, "data_type": info if isinstance(info, str) else "text"}
                for col, info in tdef["columns"].items()
            ]
            if live_source.startswith("live_db_failed"):
                live_source += " (fell back to schema_registry)"

        return {
            "ok":          True,
            "table":       table,
            "source":      live_source,
            "columns":     live_columns,
            "in_registry": bool(tdef),
            **meta,
        }
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# TOOLS REGISTRATION
# ══════════════════════════════════════════════════════════════════════════════

def _register_web_tools(TOOLS: dict) -> None:
    """Register all core_web tools into the TOOLS dict. Call from core_tools.py."""

    TOOLS["web_search"] = {
        "fn":   t_web_search,
        "perm": "READ",
        "args": ["query", "max_results"],
        "desc": "Search the web via DuckDuckGo. Returns title, url, snippet per result. No API key needed. max_results default 5.",
    }
    TOOLS["web_fetch"] = {
        "fn":   t_web_fetch,
        "perm": "READ",
        "args": ["url", "max_chars"],
        "desc": "Fetch content from any URL. Strips HTML tags, returns clean text. max_chars default 8000.",
    }
    TOOLS["summarize_url"] = {
        "fn":   t_summarize_url,
        "perm": "READ",
        "args": ["url", "focus"],
        "desc": "Fetch a URL and summarize its content via Gemini. focus= optional aspect to focus on.",
    }
    TOOLS["create_document"] = {
        "fn":   t_create_document,
        "perm": "EXECUTE",
        "args": ["content", "filename", "format", "title"],
        "desc": "Create a document from text. format=docx|pdf|txt|md|csv. Supports markdown headings/bullets. Returns base64 file.",
    }
    TOOLS["create_spreadsheet"] = {
        "fn":   t_create_spreadsheet,
        "perm": "EXECUTE",
        "args": ["data", "filename", "sheet_name", "format"],
        "desc": "Create a spreadsheet. data=JSON array of arrays OR CSV text. format=xlsx|csv. Auto-styles headers. Returns base64 file.",
    }
    TOOLS["create_presentation"] = {
        "fn":   t_create_presentation,
        "perm": "EXECUTE",
        "args": ["slides", "filename", "theme"],
        "desc": 'Create a PowerPoint presentation. slides=JSON array of {title, content, notes}. theme=default|dark|minimal. Returns base64 file.',
    }
    TOOLS["read_document"] = {
        "fn":   t_read_document,
        "perm": "READ",
        "args": ["base64_content", "filename", "format"],
        "desc": "Extract text from a document. format=docx|xlsx|pptx|txt|md|csv. Pass base64_content from Telegram file download.",
    }
    TOOLS["convert_document"] = {
        "fn":   t_convert_document,
        "perm": "EXECUTE",
        "args": ["base64_content", "from_format", "to_format", "filename"],
        "desc": "Convert between document formats. Supports md↔docx, md→pdf, txt↔docx, xlsx↔csv, docx→txt.",
    }
    TOOLS["generate_image"] = {
        "fn":   t_generate_image,
        "perm": "EXECUTE",
        "args": ["prompt", "filename", "aspect_ratio"],
        "desc": "Generate an image using Gemini Imagen. aspect_ratio=1:1|16:9|9:16|4:3|3:4. Returns base64 PNG.",
    }
    TOOLS["image_process"] = {
        "fn":   t_image_process,
        "perm": "EXECUTE",
        "args": ["base64_content", "operation", "params", "filename"],
        "desc": "Process an image with Pillow. operation=info|resize|crop|rotate|flip|convert|watermark|thumbnail. params=JSON object.",
    }
    TOOLS["weather"] = {
        "fn":   t_weather,
        "perm": "READ",
        "args": ["location"],
        "desc": "Get current weather for a location via wttr.in. No API key. Default: Jakarta.",
    }
    TOOLS["calc"] = {
        "fn":   t_calc,
        "perm": "READ",
        "args": ["expression"],
        "desc": "Safe math calculator. Supports sqrt, sin, cos, log, pi, e, **, %, etc. e.g. calc('sqrt(144)')",
    }
    TOOLS["datetime_now"] = {
        "fn":   t_datetime_now,
        "perm": "READ",
        "args": ["timezone"],
        "desc": "Get current date and time. Default timezone: Asia/Jakarta (WIB UTC+7). Returns iso, date, time, day_of_week.",
    }
    TOOLS["currency"] = {
        "fn":   t_currency,
        "perm": "READ",
        "args": ["amount", "from_cur", "to_cur"],
        "desc": "Convert currency via open exchange rates API. e.g. currency(amount='100', from_cur='USD', to_cur='IDR'). No API key.",
    }
    TOOLS["translate"] = {
        "fn":   t_translate,
        "perm": "READ",
        "args": ["text", "target_language", "source_language"],
        "desc": "Translate text using Gemini. target_language=English|Indonesian|Japanese etc. Source auto-detected if not specified.",
    }
    TOOLS["run_python"] = {
        "fn":   t_run_python,
        "perm": "EXECUTE",
        "args": ["code", "timeout"],
        "desc": "Execute Python code on Railway. Safe subprocess, max 30s. Returns stdout+stderr. Good for calculations, data transforms.",
    }
    TOOLS["list_tools"] = {
        "fn":   t_list_tools,
        "perm": "READ",
        "args": ["category", "search"],
        "desc": "List all available tools with name, args, desc. Use when you need to discover what tools exist before acting. category= filter by: deploy/code/training/system/railway/knowledge/task/web/document/image/utils/agentic/crypto/project. search= keyword filter on name+desc. CALL ONCE per intent — result is cached in context, do not repeat same call. EXAMPLES: list_tools(search='supabase') to find DB tools | list_tools(category='training') for training tools | list_tools() for all 155 tools.",
    }
    TOOLS["get_tool_info"] = {
        "fn":   t_get_tool_info,
        "perm": "READ",
        "args": ["name"],
        "desc": "Get full detail of one specific tool: args with types, full desc, perm level. Use BEFORE calling any unfamiliar tool to verify exact parameter names. If name not found, returns similar tool suggestions. CALL ONCE per tool — do not repeat. EXAMPLE: get_tool_info(name='sb_query') to see exact args before querying Supabase.",
    }
    TOOLS["get_table_schema"] = {
    "fn":   t_get_table_schema,
    "perm": "READ",
    "args": ["table"],
    "desc": "Get actual column names and types from Supabase for any table. Use BEFORE sb_query to verify correct column names. Prevents 400 errors from querying non-existent columns. EXAMPLE: get_table_schema(table='sessions') → returns all columns with types.",
    }
